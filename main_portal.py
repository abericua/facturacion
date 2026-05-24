# -*- coding: utf-8 -*-
import streamlit as st
import os
import json
import hashlib
import pyotp
import qrcode
import io
import base64
import pandas as pd
from datetime import datetime
import re

try:
    import pypdf
except ImportError:
    try:
        import PyPDF2 as pypdf
    except ImportError:
        pypdf = None

import requests # Para conectar con Gemma Local
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

import io

# --- CONFIGURACIÓN DE RUTAS PERSISTENTES (Railway Volume) ---
SGSP_ROOT = os.path.dirname(os.path.abspath(__file__))
PERSISTENT_DIR = "/app/data" # Mount path del volumen Railway
DATABASE_DIR = PERSISTENT_DIR if os.path.exists(PERSISTENT_DIR) else os.path.join(SGSP_ROOT, "database")

# Bootstrap del Volumen
if os.path.exists(PERSISTENT_DIR):
    import shutil
    repo_db = os.path.join(SGSP_ROOT, "database")
    if os.path.exists(repo_db):
        _cred_src = os.path.join(repo_db, "usuarios.json")
        _cred_dst = os.path.join(PERSISTENT_DIR, "usuarios.json")

        if not os.path.exists(_cred_dst) and os.path.exists(_cred_src):
            # Primera vez: copiar directamente
            shutil.copy2(_cred_src, _cred_dst)
        elif os.path.exists(_cred_dst) and os.path.exists(_cred_src):
            # Ya existe en el volumen: hacer merge preservando totp_secret
            try:
                import json as _json
                with open(_cred_src, 'r', encoding='utf-8', errors='replace') as f:
                    repo_users = _json.load(f)
                with open(_cred_dst, 'r', encoding='utf-8', errors='replace') as f:
                    vol_users = _json.load(f)
                vol_secrets = {u['usuario']: u.get('totp_secret', '') for u in vol_users}
                for u in repo_users:
                    if u['usuario'] in vol_secrets and vol_secrets[u['usuario']]:
                        u['totp_secret'] = vol_secrets[u['usuario']]
                with open(_cred_dst, 'w', encoding='utf-8') as f:
                    _json.dump(repo_users, f, ensure_ascii=False, indent=4)
            except Exception as _e:
                print(f"[Bootstrap] Merge warning: {_e}")

        # Copiar el resto SOLO si no existen en el volumen
        for item in os.listdir(repo_db):
            if item == "usuarios.json": continue
            s = os.path.join(repo_db, item)
            d = os.path.join(PERSISTENT_DIR, item)
            if not os.path.exists(d) and os.path.exists(s):
                shutil.copy2(s, d)

SALES_FILE = os.path.join(DATABASE_DIR, "VENTAS TOTALES 2026.xlsx")


_USERS_FILE = os.path.join(DATABASE_DIR, "usuarios.json")
USERS_FILE = _USERS_FILE
PRODUCTS_FILE = os.path.join(DATABASE_DIR, "productos_maestros.csv")

# --- AUDITORÃA DE SINCRONIZACIÓN AL ARRANQUE ---
def sync_master_data():
    if not os.path.exists(DATABASE_DIR):
        os.makedirs(DATABASE_DIR)
    
    # Buscar archivos en subcarpetas si no están en database
    mapa_sincro = {
        "usuarios.json": os.path.join(SGSP_ROOT, "Creador de Facturas", "usuarios.json"),
        "VENTAS TOTALES 2026.xlsx": os.path.join(SGSP_ROOT, "Creador de Facturas", "VENTAS TOTALES 2026.xlsx"),
        "productos_maestros.csv": os.path.join(SGSP_ROOT, "Calculadora de precios solpro", "productos_maestros.csv")
    }
    
    import shutil
    for filename, old_path in mapa_sincro.items():
        dest = os.path.join(DATABASE_DIR, filename)
        # Sincronización forzada para asegurar credenciales reales
        if os.path.exists(old_path):
            try:
                shutil.copy2(old_path, dest)
                print(f"Sincronizado: {filename}")
            except: pass

# # sync_master_data() # DESACTIVADO: Causaba pérdida de datos en producción

# --- SEGURIDAD 007 ---
SYSTEM_PEPPER = os.environ.get("SYSTEM_PEPPER", "SOLPRO_ULTRA_SECRET_2026_#!")

def hash_password(password):
    """Hash con pepper (versión nueva)"""
    salted_pass = password + SYSTEM_PEPPER
    return hashlib.sha256(salted_pass.encode()).hexdigest()

def hash_password_plain(password):
    """Hash sin pepper (versión original del facturador)"""
    return hashlib.sha256(password.encode()).hexdigest()

def verify_password(password, stored_hash):
    """Verifica contra ambos métodos de hash para compatibilidad total"""
    return (hash_password(password) == stored_hash or 
            hash_password_plain(password) == stored_hash)

def load_users():
    if os.path.exists(USERS_FILE):
        with open(USERS_FILE, 'r', encoding='utf-8', errors='replace') as f:
            return json.load(f)
    return []

def save_users(users):
    with open(USERS_FILE, 'w', encoding='utf-8') as f:
        json.dump(users, f, ensure_ascii=False, indent=4)

# --- PERSISTENCIA FINANCIERA (ALA FINANCIERA) ---
FINANZAS_FILE = os.path.join(DATABASE_DIR, "finanzas_pro.json")

def load_finanzas():
    if not os.path.exists(FINANZAS_FILE):
        init_data = {
            "bancos": [
                {"id": "c1", "banco": "Atlas", "nombre": "Atlas Cta Cte GS", "moneda": "GS", "saldo_inicial": 0},
                {"id": "c2", "banco": "Ueno", "nombre": "Ueno Digital GS", "moneda": "GS", "saldo_inicial": 0},
                {"id": "c3", "banco": "FIC", "nombre": "FIC Inversiones GS", "moneda": "GS", "saldo_inicial": 0}
            ],
            "egresos": [],
            "iva_records": {}
        }
        with open(FINANZAS_FILE, 'w', encoding='utf-8') as f:
            json.dump(init_data, f, ensure_ascii=False, indent=4)
        return init_data
    with open(FINANZAS_FILE, 'r', encoding='utf-8', errors='replace') as f:
        return json.load(f)

def save_finanzas(data):
    with open(FINANZAS_FILE, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

# --- PARSER BANCARIO (CONCILIACIÓN) ---
def parse_bank_pdf(file, bank_name):
    if pypdf is None:
        return None, "Librería pypdf no instalada. Ejecute: pip install pypdf"
    
    try:
        reader = pypdf.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        
        lines = text.split("\n")
        movimientos = []
        re_fecha = r"(\d{1,2})[/-](\d{1,2})[/-](\d{2,4})"
        
        for line in lines:
            fecha_match = re.search(re_fecha, line)
            if fecha_match:
                montos = re.findall(r"[\d]{1,3}(?:\.[\d]{3})*(?:,[\d]{2})?", line)
                nums = []
                for m in montos:
                    val = m.replace(".", "").replace(",", ".")
                    try:
                        fval = float(val)
                        if fval > 1000:
                            nums.append(fval)
                    except: pass
                
                if nums:
                    movimientos.append({
                        "fecha": "-".join(fecha_match.groups()[::-1]),
                        "descripcion": line[:100].strip(),
                        "monto": nums[0]
                    })
        
        return movimientos, None
    except Exception as e:
        return None, str(e)

def parse_iva_120(text):
    """Extrae datos del Formulario 120 (Paraguay) basado en lógica de Marangatu"""
    data = {}
    # Casillas clave
    def get_casilla(num, txt):
        pattern = rf"\b{num}\s+(\d[\d.]{{2,}})"
        matches = re.findall(pattern, txt)
        if matches:
            # Tomar el valor más alto encontrado para esa casilla
            vals = [float(m.replace(".","").replace(",",".")) for m in matches]
            return max(vals)
        return 0

    data['ventas_brutas'] = get_casilla(18, text)
    data['debito_fiscal'] = get_casilla(44, text) or get_casilla(24, text)
    data['credito_fiscal'] = get_casilla(45, text) or get_casilla(43, text)
    data['compras_netas'] = get_casilla(32, text) + get_casilla(33, text) + get_casilla(35, text) + get_casilla(36, text)
    
    # Periodo (Mes/Año)
    periodo_match = re.search(r"Mes\s+A[nñ]o\s+([\d][\d\s]{5,15})", text)
    if periodo_match:
        digits = periodo_match.group(1).replace(" ","").zfill(6)
        data['mes'] = digits[:2]
        data['anio'] = digits[2:6]
    
    return data

def scan_historical_archives():
    base_path = os.path.join(SGSP_ROOT, "IIVAS Y DOCUMENTOS LEGALES SOLPRO")
    if not os.path.exists(base_path): return []
    
    historical_results = []
    for root, dirs, files in os.walk(base_path):
        for file in files:
            if file.lower().endswith(".pdf"):
                path = os.path.join(root, file)
                try:
                    reader = pypdf.PdfReader(path)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                    
                    if "Formulario:120" in text or "VALOR AGREGADO" in text:
                        res = parse_iva_120(text)
                        res['archivo'] = file
                        historical_results.append(res)
                except: pass
    return historical_results

# --- MOTOR DE INTELIGENCIA ESTRATÉGICA (CONTEXTO TOTAL REAL) ---
def get_business_context():
    context = "### DATOS REALES DE SOLPRO (FUENTE DE VERDAD):\n"
    
    # Datos de Ventas Detallados
    if os.path.exists(SALES_FILE):
        try:
            df_v = pd.read_excel(SALES_FILE)
            df_v_active = df_v[df_v['DESCRIPCION'] != "ANULADA"]
            context += f"- VENTAS TOTALES ACUMULADAS: GS {df_v_active['PRECIO GS'].sum():,.0f}\n"
            context += f"- TOTAL OPERACIONES: {len(df_v_active)}\n"
            
            # Desglose por Línea (Si hay productos maestros para cruzar)
            prod_file = os.path.join(DATABASE_DIR, "productos_maestros.csv")
            if os.path.exists(prod_file):
                df_p = pd.read_csv(prod_file)
                # Intento de cruce para ver pesos de líneas
                context += "- PESO DE LÃNEAS DE NEGOCIO:\n"
                industrial = df_p[df_p['Linea'].str.contains("Industrial|Máquina|Laser", case=False, na=False)]['Nombre'].tolist()
                v_ind = df_v_active[df_v_active['DESCRIPCION'].isin(industrial)]['PRECIO GS'].sum()
                v_com = df_v_active['PRECIO GS'].sum() - v_ind
                context += f"  * Línea Industrial: GS {v_ind:,.0f}\n"
                context += f"  * Línea Comercial/Insumos: GS {v_com:,.0f}\n"
        except: context += "- (Error leyendo detalles de ventas)\n"
    
    # Datos Financieros Exactos
    fdata = load_finanzas()
    context += "- ESTADO DE CAJA Y BANCOS:\n"
    total_liquidez = 0
    for b in fdata["bancos"]:
        context += f"  * {b['nombre']}: {b['moneda']} {b['saldo_inicial']:,.0f}\n"
        total_liquidez += b['saldo_inicial']
    context += f"- LIQUIDEZ TOTAL DISPONIBLE: GS {total_liquidez:,.0f}\n"
    
    # Datos de Stock
    prod_file = os.path.join(DATABASE_DIR, "productos_maestros.csv")
    if os.path.exists(prod_file):
        df_p = pd.read_csv(prod_file)
        if 'Stock' in df_p.columns:
            criticos = df_p[df_p['Stock'] <= 3]['Nombre'].tolist()
            context += f"- PRODUCTOS EN QUIEBRE DE STOCK ({len(criticos)}): {', '.join(criticos[:5])}\n"
    
    # Datos Históricos (2024-2025)
    hist_file = os.path.join(DATABASE_DIR, "historical_summary.json")
    if os.path.exists(hist_file):
        with open(hist_file, 'r', encoding='utf-8', errors='replace') as f:
            hist = json.load(f)
            context += f"- HISTORIAL 2024-2025: Se han procesado {len(hist)} declaraciones juradas antiguas.\n"
            # Calcular promedio histórico si hay datos
            if hist:
                v_prom = sum(h.get('ventas_brutas',0) for h in hist) / len(hist)
                context += f"  * Venta Mensual Promedio Histórica: GS {v_prom:,.0f}\n"

    # IDENTIDAD CORPORATIVA (Del Membrete)
    context += "\n### IDENTIDAD CORPORATIVA (SOLPRO S.R.L.):\n"
    context += "- RAZÓN SOCIAL: SOLPRO S.R.L.\n"
    context += "- DIRECCIÓN: Tte. Rivarola, Fernando de la Mora (Zona Norte), Paraguay.\n"
    context += "- TELÉFONO: +(595) 986 210765\n"
    context += "- WEB: www.solpropy.com\n"
    context += "- PARTNERS OFICIALES: EPSON (Authorised Service Centre), TUCANO.\n"
            
    return context

def ask_gemma(prompt, system_context):
    """
    Alias de compatibilidad — delega a ask_ai_assistant().
    Mantener este nombre para no romper las llamadas existentes en el portal.
    """
    return ask_ai_assistant(prompt, system_context)


def ask_ai_assistant(prompt, system_context):
    """
    Motor de IA con switch de entorno:
      - LOCAL      → LM Studio en LAN (Gemma 4, 192.168.100.198:1234)
      - PRODUCTION → Google Gemini API (requiere GOOGLE_API_KEY en Railway)

    Setear ENV_MODE=PRODUCTION en las variables de Railway.
    """
    entorno = os.environ.get("ENV_MODE", "LOCAL")

    # ── ENTORNO LOCAL: Gemma en LAN ────────────────────────────────────────
    if entorno == "LOCAL":
        ip_local = os.environ.get("GEMMA_LAN_IP", "192.168.100.198")
        puerto   = os.environ.get("GEMMA_LAN_PORT", "1234")
        url      = f"http://{ip_local}:{puerto}/v1/chat/completions"
        try:
            response = requests.post(url, json={
                "model":       "google/gemma-4-e4b:2",
                "messages": [
                    {"role": "system",  "content": system_context},
                    {"role": "user",    "content": prompt}
                ],
                "temperature": 0.7,
                "stream":      False
            }, timeout=60)
            if response.status_code == 200:
                return response.json()['choices'][0]['message']['content']
            return f"Error del servidor local ({response.status_code}): {response.text}"
        except Exception as e:
            return (f"IA Local no disponible ({ip_local}:{puerto}): {e}. "
                    f"Verificá que LM Studio esté activo.")

    # ── ENTORNO PRODUCCIÓN: Google Gemini API ──────────────────────────────
    else:
        api_key = os.environ.get("GOOGLE_API_KEY", "")
        if not api_key:
            return ("IA desactivada en producción: GOOGLE_API_KEY no configurada. "
                    "Agregala como variable de entorno en Railway.")
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
            payload = {
                "contents": [
                    {
                        "parts": [
                            {"text": f"CONTEXTO DEL SISTEMA:\n{system_context}\n\nCONSULTA:\n{prompt}"}
                        ]
                    }
                ],
                "generationConfig": {
                    "temperature": 0.7,
                    "maxOutputTokens": 2048
                }
            }
            response = requests.post(url, json=payload, timeout=30)
            if response.status_code == 200:
                return response.json()['candidates'][0]['content']['parts'][0]['text']
            return f"Error Gemini API ({response.status_code}): {response.text}"
        except Exception as e:
            return f"Error conectando con Gemini API: {e}"

# --- GENERADORES MULTIMEDIA ---
def create_pptx_from_text(text):
    if Presentation is None: return None
    prs = Presentation()
    
    # Intentar parsear slides básicas por saltos de línea doble
    slides_content = text.split("\n\n")
    for content in slides_content[:10]: # Limitar a 10 slides
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        lines = content.split("\n")
        title_shape.text = lines[0][:50]
        if len(lines) > 1:
            body_shape.text = "\n".join(lines[1:10])
            
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

def get_ai_image_url(prompt):
    # Usando Pollinations.ai como motor de renderizado rápido y gratuito
    safe_prompt = requests.utils.quote(prompt)
    return f"https://image.pollinations.ai/prompt/{safe_prompt}?width=1024&height=1024&nologo=true"

# --- CONFIGURACIÓN DE PÃGINA ---
st.set_page_config(page_title="SGSP - Solpro Master Control", layout="wide", page_icon="ðŸ›ï¸")

# --- ESTILOS ANTIGRAVITY / BENTO ---
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;700;900&display=swap');
    
    :root {
        --solpro-gold: #D4AF37;
        --luxury-black: #000000;
        --border-color: #1a1a1a;
        --text-main: #ffffff;
    }

    .stApp { 
        background-color: var(--luxury-black);
        color: var(--text-main); 
        font-family: 'Inter', sans-serif; 
    }
    
    .stSidebar {
        background-color: #000000 !important;
        border-right: 1px solid var(--border-color);
    }

    /* Tarjetas de Lujo Minimalista */
    .bento-card {
        background: #050505;
        border: 1px solid var(--border-color);
        border-radius: 0px !important;
        padding: 4rem 2rem;
        transition: all 0.2s ease-in-out;
        text-align: center;
        height: 280px;
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        position: relative;
    }
    
    .bento-card:hover {
        border-color: #444;
        background: #0a0a0a;
    }

    .module-icon {
        font-size: 2rem;
        margin-bottom: 2rem;
        opacity: 0.5;
    }
    
    .module-title {
        color: #fff;
        font-size: 1.6rem;
        font-weight: 900 !important;
        letter-spacing: -0.05em;
        text-transform: uppercase;
        margin-bottom: 0.5rem;
    }
    
    .module-desc {
        color: #555;
        font-size: 0.75rem;
        font-weight: 400;
        letter-spacing: 0.1em;
        text-transform: uppercase;
    }
    
    /* Botones de Bloque (Bold & Square) */
    .stButton>button {
        background: #fff !important;
        color: #000 !important;
        border: none !important;
        border-radius: 0px !important;
        padding: 1rem !important;
        font-weight: 900 !important;
        font-size: 0.8rem !important;
        letter-spacing: 0.2em !important;
        text-transform: uppercase !important;
        transition: opacity 0.2s !important;
    }
    
    .stButton>button:hover {
        opacity: 0.9 !important;
    }

    /* Login Box */
    .stTextInput input {
        border-radius: 0px !important;
        border: 1px solid #1a1a1a !important;
        background: #050505 !important;
        color: #fff !important;
        padding: 0.8rem !important;
    }

    /* Títulos de Página */
    h1, h2, h3 {
        font-weight: 900 !important;
        letter-spacing: -0.05em !important;
        text-transform: uppercase !important;
    }

    hr { border-color: #1a1a1a !important; }
    </style>
""", unsafe_allow_html=True)

# --- ESTADO DE SESIÓN ---
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'user' not in st.session_state:
    st.session_state.user = None
if 'login_step' not in st.session_state:
    st.session_state.login_step = 1

# --- ESTADO DE NAVEGACIÓN ---
if "current_page" not in st.session_state:
    st.session_state.current_page = "portal"

def navigate_to(page):
    st.session_state.current_page = page
    st.rerun()

def login_screen():
    # --- LOGO EN LOGIN ---
    logo_path = os.path.join(SGSP_ROOT, "LOGO  3D FONDO NEGRO 2026.png")
    logo_html = ""
    if os.path.exists(logo_path):
        with open(logo_path, "rb") as f:
            data = base64.b64encode(f.read()).decode()
            logo_html = f'<img src="data:image/png;base64,{data}" width="180" style="margin-bottom: 20px; filter: drop-shadow(0px 0px 15px rgba(212, 175, 55, 0.4));">'

    st.markdown(logo_html, unsafe_allow_html=True)
    st.markdown("""
        <div style="text-align: center; margin-top: -20px;">
            <h1 style="font-family: 'Syne', sans-serif; font-size: 4.5rem; letter-spacing: -3px; color: #fff; margin-bottom: 0;">SOLPRO</h1>
            <p style="color: #D4AF37; letter-spacing: 0.5em; text-transform: uppercase; font-size: 0.8rem; font-weight: 300; margin-top: -10px; opacity: 0.8;">Elite Control Center</p>
        </div>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 1.2, 1])
    with col2:
        # Contenedor de Login con Glassmorphism
        with st.container(border=True):
            if st.session_state.login_step == 1:
                st.markdown("<h3 style='text-align:center; color:#fff; font-family:Syne;'>IDENTIFICACIÓN</h3>", unsafe_allow_html=True)
                with st.form("login_form", clear_on_submit=False):
                    u = st.text_input("USUARIO", placeholder="Nombre de acceso")
                    p = st.text_input("CLAVE", type="password", placeholder="â€¢â€¢â€¢â€¢â€¢â€¢â€¢â€¢")
                    if st.form_submit_button("VALIDAR IDENTIDAD", use_container_width=True):
                        users = load_users()
                        u_clean, p_clean = u.strip(), p.strip()
                        user = next((x for x in users if x['usuario'] == u_clean and verify_password(p_clean, x['password'])), None)
                        if user:
                            st.session_state.temp_user = user
                            st.session_state.login_step = 2
                            st.rerun()
                        else:
                            st.error("Acceso Denegado")
            else:
                user = st.session_state.temp_user
                
                # --- LÓGICA DE QR / SEGURIDAD ---
                if 'totp_secret' not in user or not user['totp_secret']:
                    secret = pyotp.random_base32()
                    user['totp_secret'] = secret
                    users = load_users()
                    for usr in users:
                        if usr['usuario'] == user['usuario']: usr['totp_secret'] = secret
                    save_users(users)
                    
                    st.info("Configuración 2FA Requerida")
                    uri = pyotp.totp.TOTP(secret).provisioning_uri(name=user['usuario'], issuer_name="SGSP_SOLPRO")
                    img = qrcode.make(uri)
                    buf = io.BytesIO()
                    img.save(buf)
                    st.image(buf.getvalue(), caption="Escanea con Google Authenticator", use_container_width=True)
                    st.code(f"Código manual: {secret}")
                else:
                    # Mostrar un icono de seguridad si ya está configurado
                    st.markdown("<div style='text-align:center; font-size:4rem; margin-bottom:10px;'>ðŸ›¡ï¸</div>", unsafe_allow_html=True)
                
                st.markdown(f"<p style='text-align:center; color:#fff;'>Protocolo 2FA: <b>{user['nombre']}</b></p>", unsafe_allow_html=True)
                code = st.text_input("CÓDIGO DE SEGURIDAD", max_chars=6, placeholder="000000")
                if st.button("AUTENTICAR ACCESO", use_container_width=True):
                    totp = pyotp.TOTP(user['totp_secret'])
                    if totp.verify(code):
                        st.session_state.logged_in = True
                        st.session_state.user = user
                        st.rerun()
                    else:
                        st.error("Código Inválido")
    st.stop()

# --- BLOQUEO DE ACCESO SI NO HAY LOGIN ---
if not st.session_state.get('logged_in', False) or st.session_state.get('user') is None:
    st.session_state.logged_in = False
    st.session_state.user = None
    login_screen()
    st.stop()

# --- LÓGICA DE NAVEGACIÓN ---
if st.session_state.logged_in and st.session_state.user:
    if st.session_state.current_page == "portal":
        # --- LOGO INSTITUCIONAL ---
        logo_path = os.path.join(SGSP_ROOT, "LOGO  2D FONDO NEGRO 2026.png")
        if os.path.exists(logo_path):
            with open(logo_path, "rb") as f:
                data = base64.b64encode(f.read()).decode()
                st.sidebar.markdown(
                    f"""
                    <div style="text-align: center; padding: 10px 0 30px 0;">
                        <img src="data:image/png;base64,{data}" width="130" style="filter: drop-shadow(0px 0px 10px rgba(212, 175, 55, 0.4));">
                    </div>
                    """,
                    unsafe_allow_html=True
                )
        
        # --- DASHBOARD PRINCIPAL (BENTO MENU) ---
        st.sidebar.markdown(f"👤 **Usuario:** {st.session_state.user['nombre']}")
        st.sidebar.markdown(f"ðŸŽ­ **Rol:** {st.session_state.user['rol'].upper()}")
        
        st.markdown(f"## Bienvenido, {st.session_state.user['nombre']}")
        
        # Sidebar shortcuts
        st.sidebar.markdown("---")
        if st.sidebar.button("ðŸ§® Motor de Precios"): navigate_to("calculadora")
        if st.sidebar.button("📊 Finanzas"): navigate_to("finanzas")
        if st.sidebar.button("â¬…ï¸ Cerrar Sesión"):
            st.session_state.logged_in = False
            st.rerun()

        # Grid Bento 1 + 2
        st.markdown("<div style='margin-bottom: 2rem;'></div>", unsafe_allow_html=True)
        
        c1, c2 = st.columns([1.5, 1])
        with c1:
            st.markdown("""
                <div class="bento-card">
                    <span class="module-icon">📄</span>
                    <div class="module-title">FACTURACIÓN & STOCK</div>
                    <div class="module-desc">Emisión legal, control de inventario y gestión de clientes.</div>
                </div>
            """, unsafe_allow_html=True)
            if st.button("LANZAR MÓDULO", key="btn_fact", use_container_width=True):
                navigate_to("facturador")

        with c2:
            st.markdown("""
                <div class="bento-card">
                    <span class="module-icon">ðŸ§®</span>
                    <div class="module-title">CALCULADORA</div>
                    <div class="module-desc">Precios v36.0 y ROI industrial.</div>
                </div>
            """, unsafe_allow_html=True)
            if st.button("ABRIR MOTOR", key="btn_calc", use_container_width=True):
                navigate_to("calculadora")

        st.markdown("""
            <div class="bento-card" style="height: 180px; flex-direction: row; gap: 30px; text-align: left;">
                <span class="module-icon" style="margin:0;">📊</span>
                <div>
                    <div class="module-title" style="margin:0;">FINANZAS & CONTROL LEGAL</div>
                    <div class="module-desc">Análisis 360, conciliación y archivo documental.</div>
                </div>
            </div>
        """, unsafe_allow_html=True)
        if st.button("ENTRAR AL CENTRO DE CONTROL FINANCIERO", key="btn_finanzas", use_container_width=True):
            navigate_to("finanzas")

    elif st.session_state.current_page == "facturador":
        if st.sidebar.button("â¬…ï¸ Volver al Portal"): navigate_to("portal")
        # Integración limpia mediante importación de módulo
        import importlib.util
        import sys
        facturador_dir = os.path.join(SGSP_ROOT, "Creador de Facturas")
        facturador_path = os.path.join(facturador_dir, "app.py")

        if os.path.exists(facturador_path):
            try:
                if facturador_dir not in sys.path:
                    sys.path.insert(0, facturador_dir)
                # Cargar el facturador como un módulo de Python real
                spec = importlib.util.spec_from_file_location("facturador_module", facturador_path)
                facturador = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(facturador)
                # Ejecutar la UI encapsulada
                facturador.run_facturador_app()
            except Exception as e:
                st.error(f"Error crítico al cargar el módulo Facturador: {e}")
                st.info("Esto puede deberse a un error de sintaxis en app.py o una dependencia faltante.")
        else:
            st.error(f"No se encontró el archivo: {facturador_path}")

    elif st.session_state.current_page == "finanzas":
        if st.sidebar.button("Volver al Portal"): navigate_to("portal")

        import subprocess
        import socket
        import sys
        def is_port_in_use(port):
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                return s.connect_ex(('localhost', port)) == 0

        react_dist_path = os.path.join(SGSP_ROOT, "mi-backoffice", "dist")
        port = 5173
        if not is_port_in_use(port):
            try:
                subprocess.Popen([sys.executable, "-m", "http.server", str(port), "-d", react_dist_path])
            except Exception as e:
                st.error(f"No se pudo iniciar el servidor del Backoffice: {e}")

        st.components.v1.html(f"<iframe src='http://localhost:{port}' width='100%' height='900' style='border:none; border-radius:12px; background:#07080f;'></iframe>", height=900)

    elif st.session_state.current_page == "calculadora":
        if st.sidebar.button("â¬…ï¸ Volver al Portal"): navigate_to("portal")
        st.markdown('<div class="header-container"><div class="header-title">ðŸ§® MOTOR DE PRECIOS SOLPRO ELITE</div><div style="color:var(--solpro-gold); font-weight:bold;">v35.2 INDUSTRIAL</div></div>', unsafe_allow_html=True)
        
        calc_path = os.path.join(SGSP_ROOT, "Calculadora de precios solpro", "calculadora_precios.html")
        if os.path.exists(calc_path):
            with open(calc_path, "r", encoding="utf-8", errors="replace") as f:
                html_content = f.read()
            
            # --- INTERFAZ DE SINCRONIZACIÓN (Admin Bridge) ---
            with st.expander("🔄 SINCRONIZADOR DE BASE DE DATOS MAESTRA", expanded=False):
                st.info("Utiliza este panel para persistir los cambios hechos en la calculadora (Combos, Precios, Márgenes) hacia la base de datos que usa el Facturador.")
                json_data = st.text_area("Pega aquí el 'Sync Token' de la calculadora para actualizar el sistema:", height=100)
                if st.button("ðŸš€ ACTUALIZAR TODO EL ECOSISTEMA SOLPRO"):
                    if json_data:
                        try:
                            import json
                            new_db = json.loads(json_data)
                            # Convertir JSON a CSV maestro
                            rows = []
                            for p in new_db:
                                rows.append({
                                    "Nombre": p['nombre'],
                                    "ID_Ref": p['id'],
                                    "Proveedor": p.get('prov', 'Manual'),
                                    "Linea": p['linea'],
                                    "Costo_Compra": p['costo'],
                                    "Moneda_Costo": p['moneda'],
                                    "Margen_Pct": p['margen']
                                })
                            df_new = pd.DataFrame(rows)
                            csv_path = os.path.join(DATABASE_DIR, "productos_maestros.csv")
                            df_new.to_csv(csv_path, index=False)
                            st.success("✅ Â¡SISTEMA SINCRONIZADO! Los vendedores ya pueden ver los nuevos precios y productos.")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Error de formato: {e}")
            
            # --- INYECCIÓN DINÃMICA DE DATOS ---
            prod_file = os.path.join(DATABASE_DIR, "productos_maestros.csv")
            if os.path.exists(prod_file):
                try:
                    df_p = pd.read_csv(prod_file)
                    dynamic_data = ""
                    for _, row in df_p.iterrows():
                        if pd.isna(row['ID_Ref']): continue
                        # ID,NOMBRE,LINEA,COSTO,MONEDA
                        line = f"{row['ID_Ref']},{row['Nombre']},{row['Linea']},{row['Costo_Compra']},{row['Moneda_Costo']}"
                        dynamic_data += line + "\n"
                    
                    import re
                    pattern = r"const MASTER_DATA = `.*?`;"
                    replacement = f"const MASTER_DATA = `{dynamic_data}`;"
                    html_content = re.sub(pattern, replacement, html_content, flags=re.DOTALL)
                    
                    saved_db_json = []
                    for _, row in df_p.iterrows():
                        saved_db_json.append({
                            "id": row['ID_Ref'],
                            "nombre": row['Nombre'],
                            "linea": row['Linea'],
                            "costo": row['Costo_Compra'],
                            "moneda": row['Moneda_Costo'],
                            "margen": row['Margen_Pct'],
                            "prov": row['Proveedor']
                        })
                    
                    sync_script = f"""
                    <script>
                        const originalSave = saveParams;
                        saveParams = function() {{
                            originalSave();
                            const token = JSON.stringify(db);
                            console.log('SYNC_TOKEN:', token);
                        }};
                    </script>
                    """
                    html_content = html_content.replace("</body>", sync_script + "</body>")
                    
                except Exception as e:
                    st.sidebar.error(f"Sync Error: {e}")
            
            st.components.v1.html(html_content, height=1000, scrolling=True)
        else:
            st.error("No se encontró el archivo de la calculadora.")

st.divider()
st.markdown("### Alertas del Sistema")
col_a1, col_a2 = st.columns(2)
with col_a1:
    st.success("✅ Todos los sistemas sincronizados.")
with col_a2:
    st.warning("⚠️ Recordatorio: Vencimiento de IVA en 3 dias.")
